from flask import Flask, render_template, request, send_file
from werkzeug.utils import secure_filename
import os
import io
import datetime

from PIL import Image
import PyPDF2
from docx import Document
from pptx import Presentation
import easyocr

from transformers import pipeline

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# 🔥 EasyOCR (image text extraction)
reader = easyocr.Reader(['en'], gpu=True)

# 🔥 GENERATIVE MODEL (DETAILED ANSWERS)
qa_pipeline = pipeline(
    "text2text-generation",
    model="google/flan-t5-base",
    device=0
)

context_dict = {}
qa_results = {}


def extract_text_from_file(file_path, filename):
    content = ""

    try:
        if filename.endswith(".pdf"):
            with open(file_path, 'rb') as f:
                reader_pdf = PyPDF2.PdfReader(f)
                for page in reader_pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"

        elif filename.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as f:
                content += f.read()

        elif filename.endswith(".docx"):
            doc = Document(file_path)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif filename.endswith(".pptx"):
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

        elif filename.lower().endswith((".png", ".jpg", ".jpeg")):
            result = reader.readtext(file_path, detail=0)
            text = " ".join(result)

            if not text.strip():
                text = "No readable text found in image."

            content += text + "\n"

    except Exception as e:
        print(f"Error extracting {filename}: {e}")

    return content


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        context_dict.clear()
        qa_results.clear()

        for fname in os.listdir(app.config['UPLOAD_FOLDER']):
            os.remove(os.path.join(app.config['UPLOAD_FOLDER'], fname))

        files = request.files.getlist('files')

        for file in files:
            if file and file.filename.strip():
                filename = secure_filename(file.filename)
                path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(path)

                text = extract_text_from_file(path, filename)

                context_dict[filename] = text
                qa_results[filename] = []

    return render_template(
        'Index.html',
        files=list(context_dict.keys()),
        qa_results=qa_results,
        selected_file=None
    )


@app.route('/ask', methods=['POST'])
def ask():
    file = request.form.get('selected_file')
    question = request.form.get('question')

    if not file or not question or file not in context_dict:
        return "Missing file or question", 400

    context = context_dict[file]

    # 🔥 PROMPT ENGINEERING (IMPORTANT)
    prompt = f"""
    Answer the question in a clear and detailed way using the context below.

    Context:
    {context[:2000]}

    Question:
    {question}

    Answer:
    """

    try:
        result = qa_pipeline(prompt, max_length=200, do_sample=False)
        answer = result[0]['generated_text']
    except Exception as e:
        answer = f"Error: {e}"

    qa_results[file].append({
        'question': question,
        'answer': answer
    })

    return render_template(
        'Index.html',
        files=list(context_dict.keys()),
        qa_results=qa_results,
        selected_file=file
    )


@app.route('/download')
def download():
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"QnA_{timestamp}.txt"

    content = io.StringIO()

    for fname, qas in qa_results.items():
        content.write(f"\nFILE: {fname}\n")
        content.write("-" * 50 + "\n")

        for i, qa in enumerate(qas, 1):
            content.write(f"Q{i}: {qa['question']}\n")
            content.write(f"A{i}: {qa['answer']}\n\n")

    content.seek(0)

    return send_file(
        io.BytesIO(content.read().encode()),
        as_attachment=True,
        download_name=filename,
        mimetype='text/plain'
    )


if __name__ == '__main__':
    app.run(debug=True)